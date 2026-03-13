# -*- coding: utf-8 -*-
"""
chart_builder.py — 공용 PPT 차트 빌더
테마 JSON 기반으로 슬라이드/차트 생성
"""
import json, os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.chart.data import CategoryChartData

# 기본 테마 경로
_DEFAULT_THEME = os.path.join(os.path.dirname(__file__), 'themes', 'white_clean.json')


def _load_theme(path):
    with open(path, 'r', encoding='utf-8') as f:
        return json.load(f)


class Builder:
    def __init__(self, theme_path=None):
        self.theme = _load_theme(theme_path or _DEFAULT_THEME)
        self.prs = Presentation()
        s = self.theme['slide']
        self.prs.slide_width = Inches(s['width_inches'])
        self.prs.slide_height = Inches(s['height_inches'])

    # ── 헬퍼 ──────────────────────────────────────────
    def _c(self, key):
        return RGBColor(*self.theme['colors'][key])

    def _palette(self, idx):
        p = self.theme['colors']['chart_palette']
        return RGBColor(*p[idx % len(p)])

    def _font_sz(self, key):
        return Pt(self.theme['font'][key])

    def _font_name(self):
        return self.theme['font']['name']

    def _slide(self):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = RGBColor(*self.theme['slide']['background'])
        return slide

    def _title(self, slide, text, sub=""):
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(10), Inches(0.5))
        p = tb.text_frame.paragraphs[0]
        r = p.add_run(); r.text = text
        r.font.size = self._font_sz('title_size'); r.font.color.rgb = self._c('title')
        r.font.bold = True; r.font.name = self._font_name()
        if sub:
            r2 = p.add_run(); r2.text = f"  {sub}"
            r2.font.size = self._font_sz('subtitle_size')
            r2.font.color.rgb = self._c('subtitle'); r2.font.name = self._font_name()

    def _unit(self, slide, text):
        tb = slide.shapes.add_textbox(Inches(10.5), Inches(0.35), Inches(2.5), Inches(0.35))
        p = tb.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT
        r = p.add_run(); r.text = f"(단위 : {text})"
        r.font.size = self._font_sz('page_num_size')
        r.font.color.rgb = self._c('subtitle'); r.font.name = self._font_name()

    def _pgnum(self, slide, n, total):
        tb = slide.shapes.add_textbox(Inches(12), Inches(7.05), Inches(1), Inches(0.3))
        p = tb.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT
        r = p.add_run(); r.text = f"{n}/{total}"
        r.font.size = self._font_sz('page_num_size')
        r.font.color.rgb = self._c('page_num'); r.font.name = self._font_name()

    def _hide_value_axis(self, chart):
        va = chart.value_axis
        va.has_major_gridlines = False; va.has_minor_gridlines = False; va.has_title = False
        va.tick_labels.font.size = Pt(2)
        va.tick_labels.font.color.rgb = RGBColor(*self.theme['slide']['background'])
        va.format.line.color.rgb = RGBColor(*self.theme['slide']['background'])
        va.format.line.width = Pt(0)

    # ── 차트 메서드 ────────────────────────────────────

    def add_hbar(self, labels, pcts, freqs, n, title, pg, tot, sort=True):
        slide = self._slide()
        self._title(slide, title, f"(n={n:,})")
        self._unit(slide, "%, 명")

        if sort:
            pairs = sorted(zip(labels, pcts, freqs), key=lambda x: x[1])
        else:
            pairs = list(reversed(list(zip(labels, pcts, freqs))))
        s_l, s_p, s_f = [p[0] for p in pairs], [p[1] for p in pairs], [p[2] for p in pairs]

        h = self.theme['hbar']
        cd = CategoryChartData(); cd.categories = s_l; cd.add_series(' ', s_p)
        cf = slide.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED,
            Inches(h['chart_left']), Inches(h['chart_top']),
            Inches(h['chart_width']), Inches(h['chart_height']), cd)
        chart = cf.chart; chart.has_legend = False

        ser = chart.series[0]; ser.format.fill.solid()
        ser.format.fill.fore_color.rgb = self._palette(h['bar_color_index'])

        plot = chart.plots[0]; plot.has_data_labels = True
        dl = plot.data_labels
        dl.font.size = self._font_sz('data_label_size')
        dl.font.color.rgb = self._c('label'); dl.font.name = self._font_name()
        dl.number_format = '0.0"%"'
        dl.label_position = XL_LABEL_POSITION.OUTSIDE_END

        for i in range(len(s_l)):
            pt_dl = ser.points[i].data_label; tf = pt_dl.text_frame; p = tf.paragraphs[0]
            run = p.add_run(); run.text = f"{s_p[i]:.1f}% ({int(s_f[i])})"
            run.font.size = self._font_sz('data_label_size')
            run.font.color.rgb = self._c('label'); run.font.name = self._font_name()

        cax = chart.category_axis
        cax.tick_labels.font.size = self._font_sz('label_size')
        cax.tick_labels.font.color.rgb = self._c('title'); cax.tick_labels.font.name = self._font_name()
        cax.format.line.color.rgb = self._c('axis_line')
        cax.format.line.width = Pt(h['axis_line_width'])
        self._hide_value_axis(chart); self._pgnum(slide, pg, tot)

    def add_pie(self, labels, pcts, freqs, n, title, pg, tot, donut=False):
        slide = self._slide()
        self._title(slide, title, f"(n={n:,})")
        self._unit(slide, "%, 명")

        p_cfg = self.theme['pie']
        cd = CategoryChartData(); cd.categories = labels; cd.add_series(' ', pcts)
        ct = XL_CHART_TYPE.DOUGHNUT if donut else XL_CHART_TYPE.PIE
        cf = slide.shapes.add_chart(ct,
            Inches(p_cfg['chart_left']), Inches(p_cfg['chart_top']),
            Inches(p_cfg['chart_width']), Inches(p_cfg['chart_height']), cd)
        chart = cf.chart; chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.legend.font.size = self._font_sz('legend_size')
        chart.legend.font.color.rgb = self._c('title'); chart.legend.font.name = self._font_name()

        plot = chart.plots[0]; plot.has_data_labels = True
        dl = plot.data_labels
        dl.font.size = self._font_sz('legend_size')
        dl.font.color.rgb = self._c('pie_label')
        dl.font.bold = p_cfg.get('label_bold', True); dl.font.name = self._font_name()
        dl.number_format = '0.0"%"'

        series = plot.series[0]
        for i, point in enumerate(series.points):
            point.format.fill.solid(); point.format.fill.fore_color.rgb = self._palette(i)
        for i in range(len(labels)):
            pt_dl = series.points[i].data_label; tf = pt_dl.text_frame; p = tf.paragraphs[0]
            run = p.add_run(); run.text = f"{pcts[i]:.1f}% ({int(freqs[i])})"
            run.font.size = self._font_sz('legend_size')
            run.font.color.rgb = self._c('pie_label')
            run.font.bold = True; run.font.name = self._font_name()
        self._pgnum(slide, pg, tot)

    def add_radar(self, labels, values, title, pg, tot):
        slide = self._slide()
        self._title(slide, title, "(5점 만점 기준)")
        self._unit(slide, "점")

        r_cfg = self.theme['radar']
        cd = CategoryChartData(); cd.categories = labels; cd.add_series(' ', values)
        cf = slide.shapes.add_chart(XL_CHART_TYPE.RADAR_MARKERS,
            Inches(r_cfg['chart_left']), Inches(r_cfg['chart_top']),
            Inches(r_cfg['chart_width']), Inches(r_cfg['chart_height']), cd)
        chart = cf.chart; chart.has_legend = False

        ser = chart.series[0]
        ser.format.line.color.rgb = self._palette(r_cfg['line_color_index'])
        ser.format.line.width = Pt(r_cfg['line_width'])
        ser.marker.style = 8; ser.marker.size = r_cfg['marker_size']
        ser.marker.format.fill.solid()
        ser.marker.format.fill.fore_color.rgb = self._palette(r_cfg['marker_color_index'])

        plot = chart.plots[0]; plot.has_data_labels = True
        dl = plot.data_labels
        dl.font.size = self._font_sz('data_label_size')
        dl.font.color.rgb = self._palette(r_cfg['line_color_index'])
        dl.font.bold = True; dl.font.name = self._font_name(); dl.number_format = '0.00'

        cax = chart.category_axis
        cax.tick_labels.font.size = self._font_sz('legend_size')
        cax.tick_labels.font.color.rgb = self._c('title'); cax.tick_labels.font.name = self._font_name()
        self._pgnum(slide, pg, tot)

    def add_stacked(self, items, data, univs, title, pg, tot):
        slide = self._slide()
        self._title(slide, title, f"(상위 {len(items)}개)")
        self._unit(slide, "명")

        st = self.theme['stacked']
        r_items = list(reversed(items))
        cd = CategoryChartData(); cd.categories = r_items
        for u in univs:
            cd.add_series(u, list(reversed(data[u])))

        cf = slide.shapes.add_chart(XL_CHART_TYPE.BAR_STACKED,
            Inches(st['chart_left']), Inches(st['chart_top']),
            Inches(st['chart_width']), Inches(st['chart_height']), cd)
        chart = cf.chart; chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.TOP; chart.legend.include_in_layout = False
        chart.legend.font.size = self._font_sz('data_label_size')
        chart.legend.font.color.rgb = self._c('title'); chart.legend.font.name = self._font_name()

        for si, s in enumerate(chart.series):
            s.format.fill.solid(); s.format.fill.fore_color.rgb = self._palette(si)

        plot = chart.plots[0]; plot.has_data_labels = True
        dl = plot.data_labels
        dl.font.size = self._font_sz('stacked_label_size')
        dl.font.color.rgb = self._c('pie_label'); dl.font.name = self._font_name()
        dl.number_format = st.get('label_format', '#,##0')
        dl.label_position = XL_LABEL_POSITION.CENTER

        cax = chart.category_axis
        cax.tick_labels.font.size = Pt(9)
        cax.tick_labels.font.color.rgb = self._c('title'); cax.tick_labels.font.name = self._font_name()
        cax.format.line.color.rgb = self._c('axis_line'); cax.format.line.width = Pt(0.5)
        self._hide_value_axis(chart); self._pgnum(slide, pg, tot)

    def save(self, path):
        os.makedirs(os.path.dirname(path), exist_ok=True)
        self.prs.save(path)
