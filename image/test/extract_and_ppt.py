# -*- coding: utf-8 -*-
"""
PDF에서 수식 영역을 이미지로 크롭·저장하고, 편집 가능한 PPT로 생성합니다.
수식 영역의 좌표는 이전 분석 결과를 기반으로 지정합니다.
"""

import sys
import os
sys.stdout.reconfigure(encoding='utf-8', errors='replace')

import fitz  # PyMuPDF
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── 경로 ──
PDF_PATH = r"d:\git_rk\image\test\울산 일자리 미스매치 모델_PPT test.pdf"
OUTPUT_DIR = r"d:\git_rk\image\test\formula_images"
PPT_PATH = r"d:\git_rk\image\test\울산_미스매치_수식_PPT.pptx"

os.makedirs(OUTPUT_DIR, exist_ok=True)

# PDF 좌표 기준 (72 DPI 기준) 수식 영역 정의
# page: 0-indexed, rect: (x0, y0, x1, y1) in PDF points
FORMULA_REGIONS = [
    {
        "page": 0,
        "rect": (50, 625, 545, 665),   # G_target 수식
        "title": "연간 취업 대상 자원 (G_target)",
        "section": "2.1. [공급] 연간 취업 대상 자원 확정",
        "description": "잠재적 구직 총량을 산출하는 기본 수식",
        "filename": "formula_01_G_target.png",
    },
    {
        "page": 1,
        "rect": (50, 190, 545, 245),   # S^(ind)_k 수식
        "title": "산업별 공급 인력 (S^(ind)_k)",
        "section": "2.2. [공급] 산업별 유효 공급 인력",
        "description": "울산 정주 희망자에 산업별 선호 확률을 적용한 인력 산출",
        "filename": "formula_02_S_ind.png",
    },
    {
        "page": 1,
        "rect": (50, 345, 545, 400),   # S^(occ)_j 수식
        "title": "직무별 공급 인력 (S^(occ)_j)",
        "section": "2.2. [공급] 직무별 유효 공급 인력",
        "description": "울산 정주 희망자에 직무별 선호 확률을 적용한 인력 산출",
        "filename": "formula_03_S_occ.png",
    },
    {
        "page": 2,
        "rect": (50, 78, 545, 115),    # M^(ind)_quant 수식 [산업]
        "title": "정량적 미스매치 - 산업 (M^(ind)_quant,k)",
        "section": "3.1. 정량적 미스매치 (수급차) - 산업",
        "description": "산업별 수급 과부족 규모 산출 (+: 초과공급, -: 초과수요)",
        "filename": "formula_04_M_ind_quant.png",
    },
    {
        "page": 2,
        "rect": (50, 122, 545, 160),   # M^(occ)_quant 수식 [직무]
        "title": "정량적 미스매치 - 직무 (M^(occ)_quant,j)",
        "section": "3.1. 정량적 미스매치 (수급차) - 직무",
        "description": "직무별 수급 과부족 규모 산출 (+: 초과공급, -: 초과수요)",
        "filename": "formula_05_M_occ_quant.png",
    },
    {
        "page": 2,
        "rect": (50, 228, 545, 265),   # 수급차율 보조 지표
        "title": "수급차율 보조 지표",
        "section": "3.1. 정량적 미스매치 - 보조 지표",
        "description": "수요 대비 과부족 비중을 백분율로 표시",
        "filename": "formula_06_ratio.png",
    },
    {
        "page": 2,
        "rect": (50, 328, 545, 378),   # M^(ind)_qual [산업]
        "title": "정성적 미스매치 - 산업 (M^(ind)_qual,k)",
        "section": "3.2. 정성적 미스매치 (구조적 불일치) - 산업",
        "description": "산업별 구성비 차이를 분석하여 인력 양성 방향성 점검",
        "filename": "formula_07_M_ind_qual.png",
    },
    {
        "page": 2,
        "rect": (50, 388, 545, 438),   # M^(occ)_qual [직무]
        "title": "정성적 미스매치 - 직무 (M^(occ)_qual,j)",
        "section": "3.2. 정성적 미스매치 (구조적 불일치) - 직무",
        "description": "직무별 구성비 차이를 분석하여 인력 양성 방향성 점검",
        "filename": "formula_08_M_occ_qual.png",
    },
]

# ── 1단계: PDF에서 수식 영역 크롭 저장 ──
print("=" * 60)
print(" PDF 수식 이미지 추출 및 PPT 생성")
print("=" * 60)

doc = fitz.open(PDF_PATH)
DPI = 300
ZOOM = DPI / 72

saved_images = []

print(f"\n[1단계] 수식 영역 크롭 및 이미지 저장 ({len(FORMULA_REGIONS)}개)")

for idx, region in enumerate(FORMULA_REGIONS):
    page = doc[region["page"]]
    
    # 크롭 영역 설정 (PDF points)
    clip = fitz.Rect(region["rect"])
    
    # 고해상도 렌더링
    mat = fitz.Matrix(ZOOM, ZOOM)
    pix = page.get_pixmap(matrix=mat, clip=clip, alpha=False)
    
    # PIL Image로 변환
    img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
    
    # 여백 추가 (패딩 20px)
    padding = 20
    padded = Image.new("RGB", (img.width + padding*2, img.height + padding*2), (255, 255, 255))
    padded.paste(img, (padding, padding))
    
    # 저장
    img_path = os.path.join(OUTPUT_DIR, region["filename"])
    padded.save(img_path, "PNG", dpi=(DPI, DPI))
    saved_images.append(img_path)
    
    print(f"  [{idx+1}] {region['title']}")
    print(f"      크기: {padded.width}x{padded.height}px -> {region['filename']}")

doc.close()

# ── 2단계: PPT 생성 ──
print(f"\n[2단계] 편집 가능한 PPT 생성")

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

total = len(FORMULA_REGIONS)

for idx, region in enumerate(FORMULA_REGIONS):
    img_path = saved_images[idx]
    
    # 빈 슬라이드
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 배경 (어두운 남색)
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = RGBColor(0x1B, 0x1B, 0x2F)
    
    # 상단 악센트 라인
    accent = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, Inches(0.06))
    accent.fill.solid()
    accent.fill.fore_color.rgb = RGBColor(0x4F, 0xC3, 0xF7)
    accent.line.fill.background()
    
    # 섹션 라벨
    sec_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(10), Inches(0.4))
    sec_run = sec_box.text_frame.paragraphs[0].add_run()
    sec_run.text = region["section"]
    sec_run.font.size = Pt(14)
    sec_run.font.color.rgb = RGBColor(0x4F, 0xC3, 0xF7)
    
    # 제목
    ttl_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.0), Inches(11), Inches(0.7))
    ttl_run = ttl_box.text_frame.paragraphs[0].add_run()
    ttl_run.text = region["title"]
    ttl_run.font.size = Pt(32)
    ttl_run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    ttl_run.font.bold = True
    
    # 설명
    dsc_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.7), Inches(11), Inches(0.5))
    dsc_run = dsc_box.text_frame.paragraphs[0].add_run()
    dsc_run.text = region["description"]
    dsc_run.font.size = Pt(16)
    dsc_run.font.color.rgb = RGBColor(0xB0, 0xB0, 0xB0)
    
    # 수식 이미지 배경 (밝은 영역)
    formula_bg = slide.shapes.add_shape(
        1, Inches(1.0), Inches(2.8), Inches(11.333), Inches(3.0)
    )
    formula_bg.fill.solid()
    formula_bg.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    formula_bg.line.color.rgb = RGBColor(0x3A, 0x3A, 0x5C)
    formula_bg.line.width = Pt(1)
    # 모서리 반올림
    formula_bg.shadow.inherit = False
    
    # 수식 이미지 삽입 (중앙 배치, 적절한 크기)
    img = Image.open(img_path)
    img_w_inch = img.width / DPI
    img_h_inch = img.height / DPI
    
    # 최대 크기 내에서 비율 유지
    max_w = 10.0
    max_h = 2.5
    scale = min(max_w / img_w_inch, max_h / img_h_inch, 3.0)  # 최대 3배 확대
    
    final_w = img_w_inch * scale
    final_h = img_h_inch * scale
    
    # 중앙 배치 계산
    center_x = 1.0 + (11.333 - final_w) / 2
    center_y = 2.8 + (3.0 - final_h) / 2
    
    pic = slide.shapes.add_picture(
        img_path,
        Inches(center_x), Inches(center_y),
        Inches(final_w), Inches(final_h)
    )
    
    # 슬라이드 번호
    pn = slide.shapes.add_textbox(Inches(12.0), Inches(7.0), Inches(1), Inches(0.3))
    pn_para = pn.text_frame.paragraphs[0]
    pn_para.alignment = PP_ALIGN.RIGHT
    pn_run = pn_para.add_run()
    pn_run.text = f"{idx+1} / {total}"
    pn_run.font.size = Pt(10)
    pn_run.font.color.rgb = RGBColor(0x60, 0x60, 0x80)
    
    print(f"  [{idx+1}/{total}] {region['title']} - 슬라이드 추가 완료")

prs.save(PPT_PATH)
print(f"\nPPT 저장 완료: {PPT_PATH}")
print(f"총 {total}장의 슬라이드 생성됨")
print(f"\n이미지 저장 폴더: {OUTPUT_DIR}")
print(f"저장된 이미지 목록:")
for img_path in saved_images:
    print(f"  - {os.path.basename(img_path)}")
