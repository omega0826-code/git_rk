# -*- coding: utf-8 -*-
"""
test.md 파일의 LaTeX 수식을 편집 가능한 PPT로 변환하는 스크립트.
각 수식 블록당 1슬라이드를 생성합니다.
방법: LaTeX -> MathML -> OMML 변환을 통해 PowerPoint 네이티브 수식으로 삽입
"""

import sys
import os
import re
import copy
from lxml import etree
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import latex2mathml.converter

# 콘솔 인코딩 문제 방지
sys.stdout.reconfigure(encoding='utf-8', errors='replace')

# ── OMML 변환을 위한 XSLT 로드 ──
XSLT_PATH = None
possible_paths = [
    r"C:\Program Files\Microsoft Office\root\Office16\MML2OMML.XSL",
    r"C:\Program Files (x86)\Microsoft Office\root\Office16\MML2OMML.XSL",
    r"C:\Program Files\Microsoft Office\Office16\MML2OMML.XSL",
    r"C:\Program Files (x86)\Microsoft Office\Office16\MML2OMML.XSL",
    r"C:\Program Files\Microsoft Office\root\Office15\MML2OMML.XSL",
    r"C:\Program Files (x86)\Microsoft Office\root\Office15\MML2OMML.XSL",
]
for p in possible_paths:
    if os.path.exists(p):
        XSLT_PATH = p
        break

if XSLT_PATH is None:
    print("ERROR: MML2OMML.XSL 파일을 찾을 수 없습니다.")
    sys.exit(1)

print(f"XSLT 경로: {XSLT_PATH}")

# XSLT 변환기 미리 로드
xslt_tree = etree.parse(XSLT_PATH)
xslt_transform = etree.XSLT(xslt_tree)

# ── 수식 정의 (직접 명시) ──
# test.md 파일의 수식을 순서대로 정의
FORMULAS = [
    {
        "latex": r"G_{target} = (\text{전체 취업자} - \text{창업자}) + \text{기타(잠재 구직자)}",
        "title": "연간 취업 대상 자원 (G_target)",
        "section": "2.1. [공급] 연간 취업 대상 자원 확정",
        "description": "잠재적 구직 총량을 산출하는 기본 수식",
    },
    {
        "latex": r"S^{(ind)}_k = \sum_{h=1}^{L} \left( G_{target, h} \times \alpha_{ulsan, h} \times \beta^{(ind)}_{k|ulsan, h} \right)",
        "title": "산업별 공급 인력 (S^(ind)_k)",
        "section": "2.2. [공급] 산업별 유효 공급 인력",
        "description": "울산 정주 희망자에 산업별 선호 확률을 적용한 인력 산출",
    },
    {
        "latex": r"S^{(occ)}_j = \sum_{h=1}^{L} \left( G_{target, h} \times \alpha_{ulsan, h} \times \gamma^{(occ)}_{j|ulsan, h} \right)",
        "title": "직무별 공급 인력 (S^(occ)_j)",
        "section": "2.2. [공급] 직무별 유효 공급 인력",
        "description": "울산 정주 희망자에 직무별 선호 확률을 적용한 인력 산출",
    },
    {
        "latex": r"M^{(ind)}_{quant, k} = S^{(ind)}_k - D^{(ind)}_k",
        "title": "정량적 미스매치 - 산업",
        "section": "3.1. 정량적 미스매치 (수급차) - 산업",
        "description": "산업별 수급 과부족 규모 산출 (+: 초과공급, -: 초과수요)",
    },
    {
        "latex": r"M^{(occ)}_{quant, j} = S^{(occ)}_j - D^{(occ)}_j",
        "title": "정량적 미스매치 - 직무",
        "section": "3.1. 정량적 미스매치 (수급차) - 직무",
        "description": "직무별 수급 과부족 규모 산출 (+: 초과공급, -: 초과수요)",
    },
    {
        "latex": r"\text{수급차율}(\%) = \frac{S_k - D_k}{D_k} \times 100",
        "title": "수급차율 보조 지표",
        "section": "3.1. 정량적 미스매치 (수급차) - 보조 지표",
        "description": "수요 대비 과부족 비중을 백분율로 표시",
    },
    {
        "latex": r"M^{(ind)}_{qual, k} = \frac{S^{(ind)}_k}{\sum S^{(ind)}} - \frac{D^{(ind)}_k}{\sum D^{(ind)}}",
        "title": "정성적 미스매치 - 산업",
        "section": "3.2. 정성적 미스매치 (구조적 불일치) - 산업",
        "description": "산업별 구성비 차이를 분석하여 인력 양성 방향성 점검",
    },
    {
        "latex": r"M^{(occ)}_{qual, j} = \frac{S^{(occ)}_j}{\sum S^{(occ)}} - \frac{D^{(occ)}_j}{\sum D^{(occ)}}",
        "title": "정성적 미스매치 - 직무",
        "section": "3.2. 정성적 미스매치 (구조적 불일치) - 직무",
        "description": "직무별 구성비 차이를 분석하여 인력 양성 방향성 점검",
    },
]


def preprocess_latex(latex_str):
    """latex2mathml 호환 형태로 전처리"""
    s = latex_str.strip()
    # \text{} -> \mathrm{} (latex2mathml이 \text를 지원하지 않을 수 있음)
    s = s.replace('\\text{', '\\mathrm{')
    # \quad -> 공백
    s = s.replace('\\quad', '\\;\\;')
    return s


def latex_to_omml(latex_str):
    """LaTeX -> MathML -> OMML 변환"""
    try:
        processed = preprocess_latex(latex_str)
        mathml_str = latex2mathml.converter.convert(processed)
        mathml_tree = etree.fromstring(mathml_str.encode('utf-8'))
        omml_tree = xslt_transform(mathml_tree)
        return omml_tree.getroot()
    except Exception as e:
        print(f"  [WARN] 변환 오류: {e}")
        return None


def add_omml_to_paragraph(paragraph, omml_elem):
    """OMML 요소를 paragraph XML에 삽입"""
    p_elem = paragraph._p
    elem = copy.deepcopy(omml_elem)
    
    # oMathPara 또는 oMath 태그 찾기
    tag_local = etree.QName(elem.tag).localname if isinstance(elem.tag, str) else ''
    
    if tag_local in ('oMathPara', 'oMath'):
        p_elem.append(elem)
    else:
        # 자식 중에서 oMath 요소 찾기
        for child in elem:
            child_local = etree.QName(child.tag).localname if isinstance(child.tag, str) else ''
            if child_local in ('oMathPara', 'oMath'):
                p_elem.append(copy.deepcopy(child))
                return
        # 못 찾으면 전체 추가
        p_elem.append(elem)


def create_ppt(formulas_data, output_path):
    """수식을 편집 가능한 PPT로 생성"""
    prs = Presentation()
    
    # 와이드스크린 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    total = len(formulas_data)
    
    for idx, fdata in enumerate(formulas_data):
        latex = fdata["latex"]
        title = fdata["title"]
        section = fdata["section"]
        description = fdata["description"]
        
        print(f"\n[슬라이드 {idx+1}/{total}] {title}")
        print(f"  LaTeX: {latex[:60]}...")
        
        # 빈 슬라이드
        slide_layout = prs.slide_layouts[6]  # Blank
        slide = prs.slides.add_slide(slide_layout)
        
        # 배경색 (어두운 남색)
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = RGBColor(0x1B, 0x1B, 0x2F)
        
        # ── 상단 악센트 라인 ──
        accent = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, Inches(0.06))
        accent.fill.solid()
        accent.fill.fore_color.rgb = RGBColor(0x4F, 0xC3, 0xF7)
        accent.line.fill.background()
        
        # ── 섹션 라벨 ──
        sec_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(10), Inches(0.4))
        sec_para = sec_box.text_frame.paragraphs[0]
        sec_run = sec_para.add_run()
        sec_run.text = section
        sec_run.font.size = Pt(14)
        sec_run.font.color.rgb = RGBColor(0x4F, 0xC3, 0xF7)
        
        # ── 수식 제목 ──
        ttl_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.0), Inches(11), Inches(0.7))
        ttl_para = ttl_box.text_frame.paragraphs[0]
        ttl_run = ttl_para.add_run()
        ttl_run.text = title
        ttl_run.font.size = Pt(32)
        ttl_run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        ttl_run.font.bold = True
        
        # ── 설명 ──
        dsc_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.7), Inches(11), Inches(0.5))
        dsc_para = dsc_box.text_frame.paragraphs[0]
        dsc_run = dsc_para.add_run()
        dsc_run.text = description
        dsc_run.font.size = Pt(16)
        dsc_run.font.color.rgb = RGBColor(0xB0, 0xB0, 0xB0)
        
        # ── 수식 영역 배경 박스 ──
        fbg = slide.shapes.add_shape(1, Inches(1.0), Inches(2.8), Inches(11.333), Inches(2.5))
        fbg.fill.solid()
        fbg.fill.fore_color.rgb = RGBColor(0x25, 0x25, 0x40)
        fbg.line.color.rgb = RGBColor(0x3A, 0x3A, 0x5C)
        fbg.line.width = Pt(1)
        
        # ── OMML 수식 삽입 ──
        omml_elem = latex_to_omml(latex)
        
        fbox = slide.shapes.add_textbox(Inches(1.5), Inches(3.0), Inches(10.333), Inches(2.0))
        fbox.text_frame.word_wrap = True
        fpara = fbox.text_frame.paragraphs[0]
        fpara.alignment = PP_ALIGN.CENTER
        
        if omml_elem is not None:
            add_omml_to_paragraph(fpara, omml_elem)
            print(f"  -> OMML 수식 삽입 성공")
        else:
            # 대안: LaTeX 원본 텍스트 표시
            frun = fpara.add_run()
            frun.text = latex
            frun.font.size = Pt(20)
            frun.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            frun.font.italic = True
            print(f"  -> LaTeX 텍스트로 대체")
        
        # ── 하단: LaTeX 원본 ──
        lbl = slide.shapes.add_textbox(Inches(1.0), Inches(5.8), Inches(11.333), Inches(0.3))
        lbl_para = lbl.text_frame.paragraphs[0]
        lbl_run = lbl_para.add_run()
        lbl_run.text = "LaTeX 원본:"
        lbl_run.font.size = Pt(10)
        lbl_run.font.color.rgb = RGBColor(0x70, 0x70, 0x90)
        
        ltx = slide.shapes.add_textbox(Inches(1.0), Inches(6.1), Inches(11.333), Inches(0.8))
        ltx.text_frame.word_wrap = True
        ltx_para = ltx.text_frame.paragraphs[0]
        ltx_run = ltx_para.add_run()
        ltx_run.text = f"$$ {latex} $$"
        ltx_run.font.size = Pt(10)
        ltx_run.font.color.rgb = RGBColor(0x50, 0x50, 0x70)
        ltx_run.font.name = "Consolas"
        
        # ── 슬라이드 번호 ──
        pn = slide.shapes.add_textbox(Inches(12.0), Inches(7.0), Inches(1), Inches(0.3))
        pn_para = pn.text_frame.paragraphs[0]
        pn_para.alignment = PP_ALIGN.RIGHT
        pn_run = pn_para.add_run()
        pn_run.text = f"{idx+1} / {total}"
        pn_run.font.size = Pt(10)
        pn_run.font.color.rgb = RGBColor(0x60, 0x60, 0x80)
    
    prs.save(output_path)
    print(f"\nPPT 저장 완료: {output_path}")
    print(f"총 {total}장의 슬라이드 생성됨")


if __name__ == "__main__":
    output_path = r"d:\git_rk\image\test\U-Mismatch_Model_수식.pptx"
    
    print("=" * 60)
    print(" 울산형 일자리 미스매치 분석 모델 - 수식 PPT 생성기")
    print("=" * 60)
    print(f"\n총 {len(FORMULAS)}개 수식 처리 예정")
    
    create_ppt(FORMULAS, output_path)
